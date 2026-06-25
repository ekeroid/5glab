"""
EdgeVision-Overview v2 — Ericsson-branded deck.

Structure:
  PART 1 — Platform (CAMARA + LTH 5G edge service, no application).
  PART 2 — EdgeVision application, mapped onto the platform.

Updated to match the actual running system on 2026-06-23:
  - Inference goes direct to pod NodePort (NOT through nef-shim).
  - HTTP and gRPC use the same host:port shape across both zones.
  - Edge ingress is Envoy Gateway, NOT NGINX.
  - Two zones registered: LTH (L40S) and Xerces (V100).
  - Manifest: ports 8080 http + 50051 grpc, componentName=infer.
  - Liveness initialDelaySeconds=1200 (V100 cold build budget).

Run:  uv run --with python-pptx --with Pillow presentation_v2.py
"""

import os
import subprocess
import tempfile
from lxml import etree

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ─── Paths ────────────────────────────────────────────────────────────────────
TEMPLATE_PATH = "/Users/EJOHEKE/Library/CloudStorage/Dropbox/LTH/AORTA/5G-setup-2026/apps/edgevision/Quick_guide_PPT_2025.pptx"
OUTPUT_PATH = "/Users/EJOHEKE/Library/CloudStorage/Dropbox/LTH/AORTA/5G-setup-2026/apps/edgevision/EdgeVision-Overview-v2.pptx"

# ─── Ericsson Brand Colours ──────────────────────────────────────────────────
BLUE = RGBColor(0x11, 0x74, 0xE6)
TEAL = RGBColor(0x23, 0x96, 0x9A)
PURPLE = RGBColor(0x97, 0x6C, 0xF4)
CORAL = RGBColor(0xE6, 0x5D, 0x6A)
ORANGE = RGBColor(0xE6, 0x6E, 0x19)
BLACK = RGBColor(0x00, 0x00, 0x00)
GRAY1 = RGBColor(0x24, 0x24, 0x24)
GRAY2 = RGBColor(0x76, 0x76, 0x76)
GRAY3 = RGBColor(0xA0, 0xA0, 0xA0)
GRAY4 = RGBColor(0xE0, 0xE0, 0xE0)
GRAY5 = RGBColor(0xF2, 0xF2, 0xF2)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BLUE = RGBColor(0xD6, 0xE2, 0xF5)
LIGHT_TEAL = RGBColor(0xBD, 0xE0, 0xE1)
LIGHT_PURPLE = RGBColor(0xE0, 0xD3, 0xFC)
LIGHT_CORAL = RGBColor(0xEF, 0xD4, 0xD3)

# ─── Fonts ────────────────────────────────────────────────────────────────────
HEADING_FONT = "Ericsson Hilda Light"
BODY_FONT = "Ericsson Hilda"
CODE_FONT = "Courier New"


# ─── Helpers ──────────────────────────────────────────────────────────────────
def _kill_shadow(shape):
    try:
        el = shape._element
        for spPr in el.iter(qn('a:spPr')):
            for eff in spPr.findall(qn('a:effectLst')):
                spPr.remove(eff)
            etree.SubElement(spPr, qn('a:effectLst'))
    except Exception:
        pass


def _kill_all_shadows_on_slide(slide):
    spTree = slide.shapes._spTree
    for spPr in spTree.iter(qn('a:spPr')):
        for eff in spPr.findall(qn('a:effectLst')):
            spPr.remove(eff)
        etree.SubElement(spPr, qn('a:effectLst'))


def add_textbox(slide, left, top, width, height, text,
                font_size=14, bold=False, italic=False, color=BLACK,
                font_name=BODY_FONT, alignment=PP_ALIGN.LEFT,
                line_spacing=1.3):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    lines = text.split("\n")
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.alignment = alignment
        if p.runs:
            run = p.runs[0]
        else:
            run = p.add_run()
            run.text = line
        run.font.size = Pt(font_size)
        run.font.name = font_name
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.italic = italic
        pPr = p._p.get_or_add_pPr()
        lnSpc = etree.SubElement(pPr, qn('a:lnSpc'))
        spcPct = etree.SubElement(lnSpc, qn('a:spcPct'))
        spcPct.set('val', str(int(line_spacing * 100000)))
    _kill_shadow(txBox)
    return txBox


def add_box(slide, left, top, width, height, text="",
            fill_color=None, border_color=BLUE, border_width=Pt(1.5),
            font_size=12, text_color=BLACK, bold=False, rounded=True,
            font_name=BODY_FONT):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    shape.line.color.rgb = border_color
    shape.line.width = border_width
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.name = font_name
        run.font.color.rgb = text_color
        run.font.bold = bold
    _kill_shadow(shape)
    return shape


def add_line(slide, start_x, start_y, end_x, end_y, color=BLUE, width=Pt(1.5), dashed=False):
    connector = slide.shapes.add_connector(1, start_x, start_y, end_x, end_y)
    connector.line.color.rgb = color
    connector.line.width = width
    if dashed:
        connector.line.dash_style = 2
    _kill_shadow(connector)
    return connector


def add_arrow(slide, start_x, start_y, end_x, end_y, color=BLUE, width=Pt(1.5)):
    connector = slide.shapes.add_connector(1, start_x, start_y, end_x, end_y)
    connector.line.color.rgb = color
    connector.line.width = width
    ln = connector._element.find('.//' + qn('a:ln'))
    if ln is None:
        spPr = connector._element.find('.//' + qn('a:spPr'))
        ln = etree.SubElement(spPr, qn('a:ln'))
    tailEnd = etree.SubElement(ln, qn('a:tailEnd'))
    tailEnd.set('type', 'triangle')
    tailEnd.set('w', 'med')
    tailEnd.set('len', 'med')
    _kill_shadow(connector)
    return connector


def set_title(slide, text, color=BLACK):
    if slide.shapes.title:
        title = slide.shapes.title
        title.text = text
        for p in title.text_frame.paragraphs:
            for run in p.runs:
                run.font.name = HEADING_FONT
                run.font.size = Pt(36)
                run.font.color.rgb = color
                run.font.bold = False


def set_title_with_subtitle(slide, title_text, subtitle_text):
    set_title(slide, title_text)
    # Put subtitle just below the title placeholder area.
    add_textbox(slide, Inches(0.5), Inches(1.25), Inches(12.3), Inches(0.4),
                subtitle_text, font_size=16, color=GRAY2, italic=True)


def add_speaker_notes(slide, text):
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = text


def add_section_divider(text, subtitle=""):
    """A heavy section break: full-bleed blue background, white title."""
    slide = prs.slides.add_slide(prs.slide_layouts[4])  # Title Only
    # Remove default title placeholder so we control text color
    if slide.shapes.title:
        sp = slide.shapes.title._element
        sp.getparent().remove(sp)
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Inches(0), Inches(0), Inches(13.33), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = BLUE
    bg.line.fill.background()
    _kill_shadow(bg)
    # Push to back
    sp_el = bg._element
    sp_el.getparent().remove(sp_el)
    slide.shapes._spTree.insert(2, sp_el)
    # Title
    add_textbox(slide, Inches(0.8), Inches(2.6), Inches(11.7), Inches(1.4),
                text, font_size=54, bold=False, color=WHITE,
                font_name=HEADING_FONT)
    if subtitle:
        add_textbox(slide, Inches(0.8), Inches(4.0), Inches(11.7), Inches(0.8),
                    subtitle, font_size=20, color=WHITE, italic=True,
                    font_name=BODY_FONT)
    return slide


# ─── Load template, clear example slides ──────────────────────────────────────
prs = Presentation(TEMPLATE_PATH)
while len(prs.slides) > 0:
    rId = prs.slides._sldIdLst[0].get(qn('r:id'))
    prs.part.drop_rel(rId)
    prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])


# ═══════════════════════════════════════════════════════════════════════════
#  TITLE
# ═══════════════════════════════════════════════════════════════════════════

def slide_title():
    slide = prs.slides.add_slide(prs.slide_layouts[0])  # Title Page
    title_ph = slide.placeholders[0]
    title_ph.text = "EdgeVision"
    for p in title_ph.text_frame.paragraphs:
        for run in p.runs:
            run.font.name = HEADING_FONT
            run.font.size = Pt(54)
            run.font.color.rgb = BLACK
    subtitle_ph = slide.placeholders[1]
    subtitle_ph.text = "A 5G edge-compute platform —\nand a real-world application that runs on it"
    for p in subtitle_ph.text_frame.paragraphs:
        for run in p.runs:
            run.font.name = BODY_FONT
            run.font.size = Pt(20)
            run.font.color.rgb = GRAY1
    add_speaker_notes(slide,
        "KEY MESSAGE: This talk is in two halves — first the platform, then the application.\n\n"
        "Talking points:\n"
        "- Part 1: the platform = CAMARA + 5G + Kubernetes that exposes compute at the edge.\n"
        "- Part 2: EdgeVision = the YOLO-based object-detection demo that runs on it.\n"
        "- The point of the split: the platform is generic; any container can run there.\n\n"
        "REFERENCE: CAMARA project, Linux Foundation + GSMA, since 2022.")


# ═══════════════════════════════════════════════════════════════════════════
#  PART 1 — THE PLATFORM
# ═══════════════════════════════════════════════════════════════════════════

def part1_divider():
    add_section_divider("Part 1 — The Platform",
                        "A 5G + Kubernetes edge service exposed via CAMARA APIs")


def slide_why_offload():
    slide = prs.slides.add_slide(prs.slide_layouts[4])  # Title Only
    set_title(slide, "Why offload compute to the network edge?")

    y = Inches(2.0)
    cards = [
        ("Battery & heat",
         "Heavy compute drains UE battery and heats the device.",
         BLUE),
        ("Latency",
         "Edge GPU + 5G beats hyperscaler RTT.\n~10 ms RTT vs ~50–100 ms to central cloud.",
         BLUE),
        ("Cost & access",
         "GPU access without buying a GPU. Pay-per-use, on demand.",
         BLUE),
        ("Multi-device collaboration",
         "Edge becomes a coordination point — multiple UEs share state.",
         BLUE),
    ]
    cw, gap = Inches(2.95), Inches(0.15)
    for i, (title, body, c) in enumerate(cards):
        x = Inches(0.5) + i * (cw + gap)
        add_box(slide, x, y, cw, Inches(0.7), title,
                fill_color=LIGHT_BLUE, border_color=c,
                font_size=13, bold=True, text_color=c)
        add_textbox(slide, x, y + Inches(0.8), cw, Inches(3.5),
                    body, font_size=12, color=GRAY1)

    add_textbox(slide, Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.5),
                "The network already has compute near the user. The question is: "
                "how does an app actually find it and use it?",
                font_size=14, italic=True, color=GRAY1)

    add_speaker_notes(slide,
        "KEY MESSAGE: Edge compute solves four real problems: battery, latency, GPU access, collaboration.\n\n"
        "Talking points:\n"
        "- Don't anchor on one of these — different apps care about different cards.\n"
        "- AR/VR cares about latency and battery; analytics cares about cost; cooperative driving cares about collaboration.\n"
        "- The closing line teases the platform question.")


def slide_what_is_camara():
    slide = prs.slides.add_slide(prs.slide_layouts[4])
    set_title(slide, "What is CAMARA?")

    add_textbox(slide, Inches(0.5), Inches(1.85), Inches(12.3), Inches(0.6),
                "An open-source project (Linux Foundation + GSMA) that defines standard "
                "HTTP APIs for what the network can do.",
                font_size=16, color=GRAY1, italic=True)

    y = Inches(2.8)
    pillars = [
        ("Communication services",
         "Quality on Demand, Edge Discovery, Slice info,\nDevice Location, Number Verification, …",
         BLUE),
        ("Compute & AI services",
         "Edge App Management, Endpoint Discovery,\nApplication Endpoint Registration",
         TEAL),
    ]
    cw, gap = Inches(5.9), Inches(0.5)
    for i, (h, body, c) in enumerate(pillars):
        x = Inches(0.5) + i * (cw + gap)
        add_box(slide, x, y, cw, Inches(0.8), h,
                fill_color=LIGHT_BLUE if c == BLUE else LIGHT_TEAL,
                border_color=c, font_size=16, bold=True, text_color=c)
        add_textbox(slide, x + Inches(0.2), y + Inches(0.9), cw - Inches(0.4),
                    Inches(2.0), body, font_size=13, color=GRAY1)

    add_textbox(slide, Inches(0.5), Inches(6.0), Inches(12.3), Inches(1.0),
                "This talk focuses on the right column — the compute side.\n"
                "Goal: a developer asks 'put my container near this UE' and the network does it.",
                font_size=13, color=GRAY2, italic=True)

    add_speaker_notes(slide,
        "KEY MESSAGE: CAMARA = standardized HTTP APIs to the network. Two families: communication, compute.\n\n"
        "Talking points:\n"
        "- The left family is what people usually associate with CAMARA (QoS, location).\n"
        "- The right family is younger but is the topic here.\n"
        "- Stress 'open source' — anyone can implement, anyone can call.")


def slide_camara_state():
    slide = prs.slides.add_slide(prs.slide_layouts[4])
    set_title(slide, "CAMARA Edge Cloud — specification state")

    y = Inches(2.0)
    headers = ["Sub-project", "Repo", "Version", "Maturity"]
    col_starts = [Inches(0.5), Inches(3.5), Inches(7.8), Inches(9.8)]
    col_widths = [Inches(3.0), Inches(4.3), Inches(1.9), Inches(2.5)]

    for j, (h, cx, cw) in enumerate(zip(headers, col_starts, col_widths)):
        add_box(slide, cx, y, cw, Inches(0.45), h,
                fill_color=BLUE, border_color=BLUE, font_size=11,
                bold=True, text_color=WHITE, rounded=False)

    rows = [
        ("Simple Edge Discovery", "camaraproject/SimpleEdgeDiscovery", "v2.0.1", "Stable (r2.3)"),
        ("Traffic Influence", "camaraproject/TrafficInfluence", "v0.10.0", "Active"),
        ("Edge App Management", "camaraproject/EdgeApplicationManagement", "v0.1.0-alpha", "Early draft"),
        ("Optimal Edge Discovery", "camaraproject/OptimalEdgeDiscovery", "v0.1.0", "Early draft"),
        ("App Endpoint Discovery", "camaraproject/ApplicationEndpointDiscovery", "v0.1.0", "Early draft"),
        ("App Endpoint Registration", "camaraproject/ApplicationEndpointRegistration", "v0.1.0", "Early draft"),
    ]
    for i, (name, repo, ver, mat) in enumerate(rows):
        row_y = y + Inches(0.5 + i * 0.5)
        fill = WHITE if i % 2 == 0 else GRAY5
        for j, (text, cx, cw) in enumerate(zip([name, repo, ver, mat], col_starts, col_widths)):
            color = GRAY1
            if "alpha" in text or "Early" in text:
                color = CORAL
            elif "Stable" in text:
                color = TEAL
            add_box(slide, cx, row_y, cw, Inches(0.45), text,
                    fill_color=fill, border_color=GRAY4, font_size=10,
                    text_color=color, rounded=False,
                    font_name=CODE_FONT if j == 1 else BODY_FONT)

    add_textbox(slide, Inches(0.5), Inches(5.8), Inches(12.0), Inches(1.0),
                "Edge App Management is still alpha — schemas are not final.\n"
                "We follow the lifecycle semantics; manifest details can be re-aligned when v1.0 lands.",
                font_size=12, color=GRAY2, italic=True)

    add_speaker_notes(slide,
        "KEY MESSAGE: Only Edge Discovery is stable. App Management is alpha.\n\n"
        "Talking points:\n"
        "- The whole Edge Cloud umbrella was split into 6 sub-projects in April 2024.\n"
        "- Stable = SimpleEdgeDiscovery only. Everything else is moving.\n"
        "- This is why our shim diverges on manifest shape — there's nothing solid to align to yet.\n\n"
        "REFERENCE: github.com/camaraproject")


def slide_platform_arch_generic():
    """The platform diagram WITHOUT EdgeVision — just the service shape."""
    slide = prs.slides.add_slide(prs.slide_layouts[4])
    set_title(slide, "The LTH 5G edge-compute platform")
    add_textbox(slide, Inches(0.5), Inches(1.25), Inches(12.3), Inches(0.4),
                "Generic. Any container can run. Application not yet introduced.",
                font_size=14, italic=True, color=GRAY2)

    y_start = Inches(1.9)
    layers = [
        ("Client / UE", Inches(0.5), y_start, Inches(3.5), Inches(5.0), LIGHT_BLUE, BLUE),
        ("5G SA Network",   Inches(4.3), y_start, Inches(4.0), Inches(5.0), GRAY5,      GRAY3),
        ("Edge Cloud (Kubernetes)", Inches(8.6), y_start, Inches(4.2), Inches(5.0), LIGHT_TEAL, TEAL),
    ]
    for label, lx, ly, lw, lh, fill, border in layers:
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, lx, ly, lw, lh)
        bg.fill.solid(); bg.fill.fore_color.rgb = fill
        bg.line.color.rgb = GRAY4; bg.line.width = Pt(0.5)
        _kill_shadow(bg)
        add_textbox(slide, lx + Inches(0.1), ly + Inches(0.08), lw - Inches(0.2), Inches(0.4),
                    label, font_size=12, bold=True, color=GRAY2)

    # Client side
    cx = Inches(0.8)
    add_box(slide, cx, y_start + Inches(0.8), Inches(2.8), Inches(0.7),
            "Your application", fill_color=WHITE, border_color=BLUE,
            font_size=13, bold=True, text_color=BLUE)
    add_box(slide, cx, y_start + Inches(1.8), Inches(2.8), Inches(0.7),
            "CAMARA client (HTTP)", fill_color=LIGHT_BLUE, border_color=BLUE,
            font_size=12, bold=True, text_color=BLUE)
    add_box(slide, cx, y_start + Inches(2.8), Inches(2.8), Inches(0.7),
            "Data-plane client\n(any TCP / gRPC / HTTP)",
            fill_color=WHITE, border_color=BLUE, font_size=11, text_color=BLUE)

    # 5G core (kept abstract; this is the platform view, not the radio detail)
    nx = Inches(4.6)
    add_box(slide, nx, y_start + Inches(1.0), Inches(3.4), Inches(0.7),
            "gNB + 5G SA Core (Open5GS)", fill_color=WHITE, border_color=GRAY2,
            font_size=12, text_color=GRAY1)
    add_box(slide, nx, y_start + Inches(2.2), Inches(3.4), Inches(0.9),
            "NEF-shim\n(CAMARA API → k8s, multi-cluster)",
            fill_color=LIGHT_BLUE, border_color=BLUE,
            font_size=13, bold=True, text_color=BLUE)

    # Edge side — generic container, no Triton mention
    ex = Inches(8.9)
    add_box(slide, ex, y_start + Inches(0.7), Inches(3.6), Inches(0.7),
            "Envoy Gateway", fill_color=WHITE, border_color=TEAL,
            font_size=12, text_color=TEAL)
    add_box(slide, ex, y_start + Inches(1.7), Inches(3.6), Inches(0.7),
            "NodePort Services\n(one port per tenant interface)",
            fill_color=WHITE, border_color=TEAL, font_size=11, text_color=TEAL)
    add_box(slide, ex, y_start + Inches(2.7), Inches(3.6), Inches(1.0),
            "Tenant Pod\n(any GPU/CPU container)",
            fill_color=LIGHT_TEAL, border_color=TEAL,
            font_size=13, bold=True, text_color=TEAL)
    add_box(slide, ex, y_start + Inches(3.95), Inches(3.6), Inches(0.7),
            "NVIDIA L40S / V100",
            fill_color=WHITE, border_color=TEAL, font_size=11, text_color=TEAL)

    # Arrows
    # CAMARA control plane
    add_arrow(slide, Inches(3.6), y_start + Inches(2.15), Inches(4.6), y_start + Inches(2.6),
              color=BLUE, width=Pt(2))
    add_textbox(slide, Inches(3.55), y_start + Inches(2.25), Inches(1.4), Inches(0.3),
                "CAMARA", font_size=10, italic=True, color=BLUE)
    # NEF → k8s API (control)
    add_arrow(slide, Inches(8.0), y_start + Inches(2.6), Inches(8.9), y_start + Inches(3.2),
              color=GRAY2, width=Pt(1.5))
    add_textbox(slide, Inches(8.0), y_start + Inches(2.3), Inches(1.0), Inches(0.3),
                "k8s API", font_size=10, italic=True, color=GRAY2)
    # Data plane: client → pod, bypassing NEF
    add_arrow(slide, Inches(3.6), y_start + Inches(3.15), Inches(8.9), y_start + Inches(3.15),
              color=TEAL, width=Pt(2))
    add_textbox(slide, Inches(5.5), y_start + Inches(2.8), Inches(3.0), Inches(0.3),
                "data plane (direct to NodePort)", font_size=11,
                italic=True, color=TEAL, alignment=PP_ALIGN.CENTER)

    add_speaker_notes(slide,
        "KEY MESSAGE: The platform exposes three things to any application: a control-plane API (CAMARA), a place to run a container, and a way to reach it.\n\n"
        "Talking points:\n"
        "- This slide is intentionally generic. No YOLO, no Triton, no model — just the service shape.\n"
        "- Blue arrow = control plane (the CAMARA API). Teal arrow = data plane (your app's traffic).\n"
        "- The data plane DOES NOT pass through the NEF-shim. NEF returns host:port and steps aside.\n"
        "- 'Tenant Pod' is whatever the developer brings — we'll see what we put there in Part 2.\n\n"
        "Colour meaning: Blue = control / CAMARA, Teal = data plane / edge compute, Grey = transport.")


def slide_lifecycle_generic():
    slide = prs.slides.add_slide(prs.slide_layouts[4])
    set_title(slide, "CAMARA lifecycle — 4 steps, any application")

    y = Inches(2.0)
    apis = [
        ("1. Discover zones", "GET /edge-cloud-zones",
         "Which edge zones are\nreachable from this UE.", BLUE),
        ("2. Register & instantiate", "POST /apps  →  POST /app-instances",
         "Manifest in, k8s Deployment\n+ Service out.", BLUE),
        ("3. Endpoint discovery", "GET /endpoints",
         "host:port for each port\nin the manifest.", BLUE),
    ]
    for i, (title, endpoint, desc, color) in enumerate(apis):
        x = Inches(0.7 + i * 4.2)
        add_box(slide, x, y, Inches(3.8), Inches(0.9), title,
                fill_color=LIGHT_BLUE, border_color=color,
                font_size=14, bold=True, text_color=color)
        add_textbox(slide, x + Inches(0.2), y + Inches(1.0), Inches(3.6), Inches(0.4),
                    endpoint, font_size=10, font_name=CODE_FONT, color=GRAY1)
        add_textbox(slide, x + Inches(0.2), y + Inches(1.45), Inches(3.6), Inches(0.7),
                    desc, font_size=12, color=GRAY1)
        if i < 2:
            add_arrow(slide, x + Inches(3.9), y + Inches(0.45),
                      x + Inches(4.1), y + Inches(0.45), color=color, width=Pt(2))

    # Teardown
    ty = Inches(4.7)
    add_box(slide, Inches(4.9), ty, Inches(3.8), Inches(0.7),
            "4. Teardown", fill_color=WHITE, border_color=CORAL,
            font_size=14, bold=True, text_color=CORAL)
    add_textbox(slide, Inches(5.1), ty + Inches(0.8), Inches(3.4), Inches(0.4),
                "DELETE /app-instances/{id}", font_size=11,
                font_name=CODE_FONT, color=GRAY1)
    add_textbox(slide, Inches(5.1), ty + Inches(1.2), Inches(3.4), Inches(0.5),
                "k8s Deployment + Service removed.", font_size=12, color=GRAY1)
    add_arrow(slide, Inches(6.8), y + Inches(2.2), Inches(6.8), ty, color=CORAL, width=Pt(1.5))

    add_speaker_notes(slide,
        "KEY MESSAGE: Four HTTP calls cover the full lifetime of an edge workload.\n\n"
        "Talking points:\n"
        "- Step 1: discovery is geography + capabilities (GPU model, CPU cores, memory).\n"
        "- Step 2: register the manifest, then create an instance on a chosen zone. Two calls so the same manifest can be instantiated repeatedly.\n"
        "- Step 3: returns host:port. The client now bypasses CAMARA entirely.\n"
        "- Step 4: when done, DELETE. k8s resources go away; manifest stays registered if you want to redeploy.\n\n"
        "Colour meaning: Blue = lifecycle steps, Coral = teardown / destructive.")


def slide_camara_api_details():
    slide = prs.slides.add_slide(prs.slide_layouts[4])
    set_title(slide, "CAMARA API — request and response shapes")

    y = Inches(1.85)
    # Headers
    col_x = [Inches(0.5), Inches(4.1), Inches(8.8)]
    col_w = [Inches(3.5), Inches(4.6), Inches(4.0)]
    for j, (h, cx, cw) in enumerate(zip(
            ["Method + path", "Request body", "Response body"], col_x, col_w)):
        add_box(slide, cx, y, cw, Inches(0.4), h,
                fill_color=BLUE, border_color=BLUE,
                font_size=11, bold=True, text_color=WHITE, rounded=False)

    rows = [
        # (verb, path, request, response)
        ("GET",
         "/simple-edge-discovery/v0/\n  edge-cloud-zones?device-ip=…",
         "—",
         "{\n  \"edgeCloudZones\": [\n    {\"edgeCloudZoneId\": …,\n     \"capabilities\": {\n       \"gpuModel\": \"L40S\",\n       \"gpuCount\": 2, …}}\n  ]\n}"),
        ("POST",
         "/edge-app-management/v0/\n  apps",
         "{ \"appName\": …,\n  \"containerSpec\": {\n    \"imageName\": …,\n    \"imageTag\": …,\n    \"readinessProbe\": {…},\n    \"livenessProbe\":  {…}\n  },\n  \"requiredResources\":\n    {\"cpu\":2,\"memory\":8192,\n     \"gpu\":1},\n  \"componentSpec\": [{\n    \"networkInterfaces\":[…]\n  }]\n}",
         "{ \"appId\":\n    \"4c3ba6ea-…\" }"),
        ("POST",
         "/edge-app-management/v0/\n  app-instances",
         "{ \"appId\": \"4c3ba6ea-…\",\n  \"edgeCloudZoneId\":\n    \"lth-5glab-gpu-zone\" }",
         "{ \"appInstanceId\":\n    \"d6feafb4-…\",\n  \"status\":\n    \"instantiating\" }"),
        ("GET",
         "/edge-app-management/v0/\n  app-instances/{id}",
         "—",
         "{ \"status\":\n    \"instantiating | ready | failed\",\n  \"phase\":\n    \"scheduling | pulling |\n     starting | running |\n     ready | failed\",\n  \"message\": \"…\" }"),
        ("GET",
         "/application-endpoint-discovery/\n  v0/endpoints?appInstanceId=…",
         "—",
         "{ \"endpoints\": [\n    {\"name\":\"http\",\n     \"url\":\"http://host:32756\"},\n    {\"name\":\"grpc\",\n     \"url\":\"host:30661\"}\n  ]}"),
    ]

    ry = y + Inches(0.45)
    row_h = Inches(1.05)
    for i, (verb, path, req, resp) in enumerate(rows):
        fill = WHITE if i % 2 == 0 else GRAY5

        # Method + path cell
        cell = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                       col_x[0], ry, col_w[0], row_h)
        cell.fill.solid(); cell.fill.fore_color.rgb = fill
        cell.line.color.rgb = GRAY4; cell.line.width = Pt(0.5)
        _kill_shadow(cell)
        # Verb badge
        verb_color = TEAL if verb == "GET" else (BLUE if verb == "POST" else CORAL)
        add_box(slide, col_x[0] + Inches(0.1), ry + Inches(0.1),
                Inches(0.6), Inches(0.3), verb,
                fill_color=verb_color, border_color=verb_color,
                font_size=10, bold=True, text_color=WHITE, rounded=False)
        add_textbox(slide, col_x[0] + Inches(0.8), ry + Inches(0.08),
                    col_w[0] - Inches(0.9), row_h - Inches(0.1),
                    path, font_size=9, font_name=CODE_FONT, color=GRAY1,
                    line_spacing=1.15)

        # Request cell
        cell = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                       col_x[1], ry, col_w[1], row_h)
        cell.fill.solid(); cell.fill.fore_color.rgb = fill
        cell.line.color.rgb = GRAY4; cell.line.width = Pt(0.5)
        _kill_shadow(cell)
        add_textbox(slide, col_x[1] + Inches(0.1), ry + Inches(0.05),
                    col_w[1] - Inches(0.2), row_h - Inches(0.1),
                    req, font_size=8, font_name=CODE_FONT,
                    color=GRAY1 if req != "—" else GRAY3, line_spacing=1.1)

        # Response cell
        cell = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                       col_x[2], ry, col_w[2], row_h)
        cell.fill.solid(); cell.fill.fore_color.rgb = fill
        cell.line.color.rgb = GRAY4; cell.line.width = Pt(0.5)
        _kill_shadow(cell)
        add_textbox(slide, col_x[2] + Inches(0.1), ry + Inches(0.05),
                    col_w[2] - Inches(0.2), row_h - Inches(0.1),
                    resp, font_size=8, font_name=CODE_FONT,
                    color=GRAY1, line_spacing=1.1)

        ry += row_h

    # Footer with DELETE
    fy = ry + Inches(0.1)
    add_box(slide, Inches(0.5), fy, Inches(0.6), Inches(0.3), "DELETE",
            fill_color=CORAL, border_color=CORAL,
            font_size=10, bold=True, text_color=WHITE, rounded=False)
    add_textbox(slide, Inches(1.2), fy + Inches(0.02),
                Inches(11.5), Inches(0.4),
                "/edge-app-management/v0/app-instances/{id}   →   204 No Content",
                font_size=11, font_name=CODE_FONT, color=GRAY1)

    add_speaker_notes(slide,
        "KEY MESSAGE: This is the whole API. Five endpoints plus DELETE.\n\n"
        "Talking points:\n"
        "- GETs return JSON, no body needed. POSTs send JSON.\n"
        "- Two POSTs, deliberately: register manifest separately from instantiating it. Same manifest can be re-deployed without re-sending the body.\n"
        "- The polled status uses two fields: 'status' (terminal-ish) and 'phase' (fine-grained k8s phase).\n"
        "- Endpoint discovery is what the client really wants. After this, no more CAMARA calls until teardown.\n"
        "- DELETE returns 204 No Content. Idempotent: deleting a non-existent instance returns 404 but the cluster state is the same.\n\n"
        "Colour meaning: Teal verb = GET (safe), Blue verb = POST (creates), Coral verb = DELETE (destroys).")


def slide_camara_sequence_generic():
    slide = prs.slides.add_slide(prs.slide_layouts[4])
    set_title(slide, "Using the platform — generic sequence")
    add_textbox(slide, Inches(0.5), Inches(1.25), Inches(12.3), Inches(0.4),
                "Any client, any container. No EdgeVision specifics on this slide.",
                font_size=14, italic=True, color=GRAY2)

    y = Inches(1.85)
    actors = [
        ("Client",     Inches(1.3)),
        ("NEF-shim",   Inches(4.5)),
        ("K8s API",    Inches(7.7)),
        ("App Pod",    Inches(10.7)),
    ]
    for name, x in actors:
        c = LIGHT_BLUE if name in ("Client", "NEF-shim") else LIGHT_TEAL
        b = BLUE if name in ("Client", "NEF-shim") else TEAL
        add_box(slide, x - Inches(0.7), y, Inches(1.5), Inches(0.5),
                name, fill_color=c, border_color=b,
                font_size=12, bold=True, text_color=b)
        add_line(slide, x + Inches(0.05), y + Inches(0.5),
                 x + Inches(0.05), y + Inches(5.0),
                 color=GRAY3, width=Pt(1), dashed=True)

    messages = [
        (0, 1, "GET /edge-cloud-zones",                Inches(0.7), False),
        (1, 0, "200  zones[]",                          Inches(0.95), True),
        (0, 1, "POST /apps   (manifest)",               Inches(1.25), False),
        (1, 0, "200  { appId }",                        Inches(1.5), True),
        (0, 1, "POST /app-instances   { appId, zoneId }", Inches(1.8), False),
        (1, 2, "create Deployment + NodePort Service",  Inches(2.05), False),
        (1, 0, "202  { appInstanceId, status=… }",      Inches(2.3), True),
        (2, 3, "schedule + image pull + start",         Inches(2.55), False),
        (3, 3, "container init   (depends on app)",     Inches(2.85), False),
        (0, 1, "GET /app-instances/{id}   (poll every 2 s)", Inches(3.2), False),
        (1, 2, "watch pod status",                       Inches(3.45), False),
        (3, 2, "readinessProbe passes",                  Inches(3.7), False),
        (1, 0, "200  { status: ready, phase: ready }",   Inches(3.95), True),
        (0, 1, "GET /endpoints?appInstanceId=…",         Inches(4.25), False),
        (1, 0, "200  [{name, url=host:port}, …]",        Inches(4.5), True),
        (0, 3, "data plane — direct to NodePort",        Inches(4.8), False),
    ]
    for src, dst, label, dy, is_response in messages:
        sx = actors[src][1] + Inches(0.05)
        dx = actors[dst][1] + Inches(0.05)
        msg_y = y + dy
        # Pick colour per role: response = teal, request = blue, data plane = orange
        if "data plane" in label:
            color = ORANGE; width = Pt(2.5)
        elif is_response:
            color = TEAL; width = Pt(1.2)
        else:
            color = BLUE; width = Pt(1.5)
        if src == dst:
            add_arrow(slide, sx + Inches(0.3), msg_y, sx + Inches(0.3),
                      msg_y + Inches(0.25), color=GRAY2, width=Pt(1.2))
            add_textbox(slide, sx + Inches(0.4), msg_y - Inches(0.05),
                        Inches(3.0), Inches(0.25), label,
                        font_size=10, italic=True, color=GRAY2)
        else:
            add_arrow(slide, sx, msg_y, dx, msg_y, color=color, width=width)
            mid_x = min(sx, dx) + Inches(0.1)
            add_textbox(slide, mid_x, msg_y - Inches(0.22),
                        Inches(4.0), Inches(0.25), label,
                        font_size=10, color=GRAY1, font_name=CODE_FONT)

    # Legend
    ly = Inches(6.5)
    add_box(slide, Inches(0.5), ly, Inches(0.3), Inches(0.15), "",
            fill_color=BLUE, border_color=BLUE, rounded=False)
    add_textbox(slide, Inches(0.9), ly - Inches(0.05), Inches(2.5), Inches(0.3),
                "request", font_size=10, color=GRAY1)
    add_box(slide, Inches(2.7), ly, Inches(0.3), Inches(0.15), "",
            fill_color=TEAL, border_color=TEAL, rounded=False)
    add_textbox(slide, Inches(3.1), ly - Inches(0.05), Inches(2.5), Inches(0.3),
                "response", font_size=10, color=GRAY1)
    add_box(slide, Inches(5.0), ly, Inches(0.3), Inches(0.15), "",
            fill_color=ORANGE, border_color=ORANGE, rounded=False)
    add_textbox(slide, Inches(5.4), ly - Inches(0.05), Inches(7.0), Inches(0.3),
                "data plane (bypasses NEF-shim — direct to NodePort)",
                font_size=10, color=GRAY1)

    add_speaker_notes(slide,
        "KEY MESSAGE: One round-trip of CAMARA setup, then the client talks straight to the pod.\n\n"
        "Talking points:\n"
        "- Counting the CAMARA calls: 5 requests for the entire lifecycle (zones, apps, app-instances, polled status, endpoints). Then DELETE at quit.\n"
        "- The polling loop is intentional — CAMARA doesn't notify you when a pod is ready. Caller polls.\n"
        "- The orange arrow is the load-bearing point of the entire deck: once endpoints are returned, the platform steps out of the way.\n"
        "- The 'container init (depends on app)' is the only app-specific bit on this slide. For a small static-site container that's a few seconds; for a TRT-engine-building Triton it's many minutes.\n\n"
        "Colour meaning: Blue = CAMARA request, Teal = CAMARA response, Orange = data plane.")


def slide_camara_curl():
    slide = prs.slides.add_slide(prs.slide_layouts[4])
    set_title(slide, "It's just HTTP — drive it from anything")

    y = Inches(1.85)
    # Code block on the left
    code = (
        "API=http://camara.5glab.control.lth.se\n"
        "\n"
        "# 1. Discover edge zones available to this UE\n"
        "curl -s $API/simple-edge-discovery/v0/edge-cloud-zones \\\n"
        "  --get --data-urlencode device-ip=$(my_ip)\n"
        "\n"
        "# 2. Register an app manifest\n"
        "APP_ID=$(curl -s -X POST $API/edge-app-management/v0/apps \\\n"
        "  -H 'Content-Type: application/json' \\\n"
        "  --data @manifest.json  | jq -r .appId)\n"
        "\n"
        "# 3. Instantiate it on a zone\n"
        "INST=$(curl -s -X POST $API/edge-app-management/v0/app-instances \\\n"
        "  -H 'Content-Type: application/json' \\\n"
        "  -d \"{\\\"appId\\\":\\\"$APP_ID\\\",\n"
        "       \\\"edgeCloudZoneId\\\":\\\"lth-5glab-gpu-zone\\\"}\" \\\n"
        "  | jq -r .appInstanceId)\n"
        "\n"
        "# 4. Poll until ready\n"
        "while ! curl -s $API/edge-app-management/v0/app-instances/$INST \\\n"
        "         | jq -e '.status == \"ready\"' > /dev/null;\n"
        "do sleep 2; done\n"
        "\n"
        "# 5. Discover endpoints and start talking to the pod\n"
        "curl -s \"$API/application-endpoint-discovery/v0/endpoints\\\n"
        "?appInstanceId=$INST\" | jq .endpoints\n"
        "\n"
        "# (then: your app's traffic — gRPC, HTTP, raw TCP — straight to host:port)\n"
        "\n"
        "# 6. Teardown\n"
        "curl -s -X DELETE \\\n"
        "  $API/edge-app-management/v0/app-instances/$INST\n"
    )
    add_box(slide, Inches(0.5), y, Inches(8.6), Inches(5.2), "",
            fill_color=GRAY5, border_color=GRAY3, rounded=False)
    add_textbox(slide, Inches(0.65), y + Inches(0.1),
                Inches(8.3), Inches(5.0), code,
                font_size=9.5, font_name=CODE_FONT, color=GRAY1,
                line_spacing=1.15)

    # Right column annotations
    ax = Inches(9.4)
    add_textbox(slide, ax, y, Inches(3.5), Inches(0.4),
                "Any client, any language", font_size=14, bold=True, color=BLUE)
    notes = [
        ("Stock HTTP + JSON",
         "Plain HTTP. curl, requests, axios, anything."),
        ("No SDK required",
         "Spec is on /openapi.json — generate a client if you want."),
        ("Tenancy is by source IP",
         "Same calling IP = same tenant slug. No tokens to manage in the lab."),
        ("Reference client",
         "edge-demo/edge_demo.py — full lifecycle in ~140 lines of Python."),
    ]
    ny = y + Inches(0.6)
    for label, body in notes:
        add_textbox(slide, ax, ny, Inches(3.8), Inches(0.35),
                    label, font_size=12, bold=True, color=TEAL)
        add_textbox(slide, ax, ny + Inches(0.3), Inches(3.8), Inches(1.0),
                    body, font_size=11, color=GRAY1)
        ny += Inches(1.1)

    add_speaker_notes(slide,
        "KEY MESSAGE: The API is plain HTTP. You can drive it from bash.\n\n"
        "Talking points:\n"
        "- This is a complete lifecycle. Cut/paste, edit the manifest, you have your own edge container running.\n"
        "- Step 2 separation (register → instantiate) means the same manifest can be redeployed without resending the body — useful for scripted demos.\n"
        "- The polling loop is what every client does. Could be replaced with WebSocket/SSE later.\n"
        "- Point at edge-demo/edge_demo.py — that's a real, working Python client with no abstractions. Good code to read.\n\n"
        "REFERENCE: OpenAPI spec at $API/openapi.json. edge-demo/README.md documents each step in detail.")


def slide_nef_shim_role():
    slide = prs.slides.add_slide(prs.slide_layouts[4])
    set_title(slide, "What the NEF-shim does")

    y = Inches(2.0)
    add_textbox(slide, Inches(0.5), y, Inches(2.7), Inches(0.4),
                "Incoming CAMARA calls", font_size=14, bold=True, color=BLUE)
    calls = [
        "GET  /simple-edge-discovery/v0/edge-cloud-zones",
        "POST /edge-app-management/v0/apps",
        "POST /edge-app-management/v0/app-instances",
        "GET  /edge-app-management/v0/app-instances/{id}",
        "DELETE /edge-app-management/v0/app-instances/{id}",
        "GET  /application-endpoint-discovery/v0/endpoints",
    ]
    for i, c in enumerate(calls):
        cy = y + Inches(0.55 + i * 0.5)
        add_box(slide, Inches(0.5), cy, Inches(4.3), Inches(0.42),
                c, fill_color=WHITE, border_color=BLUE, font_size=9.5,
                font_name=CODE_FONT, text_color=GRAY1, rounded=False)

    # Center shim
    sx, sy = Inches(5.4), y + Inches(0.8)
    add_box(slide, sx, sy, Inches(3.0), Inches(2.6),
            "NEF-shim\n(FastAPI, multi-cluster)",
            fill_color=LIGHT_BLUE, border_color=BLUE,
            font_size=16, bold=True, text_color=BLUE)
    add_textbox(slide, sx, sy + Inches(2.7), Inches(3.0), Inches(0.6),
                "Tenant slug:\nt-{md5(source_ip)[:8]}",
                font_size=11, font_name=CODE_FONT, color=GRAY2,
                alignment=PP_ALIGN.CENTER)

    # Right column: ops
    add_textbox(slide, Inches(8.9), y, Inches(3.5), Inches(0.4),
                "Kubernetes operations (per zone)", font_size=14, bold=True, color=TEAL)
    ops = [
        "list nodes with nvidia.com/gpu.present",
        "create Deployment t-{slug}-app",
        "create Service t-{slug}-app (NodePort)",
        "watch pod status / conditions",
        "delete Deployment + Service",
        "lookup nodePort + ext_host per zone",
    ]
    for i, o in enumerate(ops):
        cy = y + Inches(0.55 + i * 0.5)
        add_box(slide, Inches(8.9), cy, Inches(3.9), Inches(0.42),
                o, fill_color=WHITE, border_color=TEAL, font_size=10,
                font_name=CODE_FONT, text_color=GRAY1, rounded=False)

    # Arrows
    for i in range(6):
        cy = y + Inches(0.76 + i * 0.5)
        add_arrow(slide, Inches(4.8), cy, Inches(5.4), cy, color=BLUE, width=Pt(1.3))
        add_arrow(slide, Inches(8.4), cy, Inches(8.9), cy, color=TEAL, width=Pt(1.3))

    add_textbox(slide, Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.4),
                "Important: the NEF-shim is control plane only. Application traffic does not pass through it.",
                font_size=13, italic=True, color=CORAL)

    add_speaker_notes(slide,
        "KEY MESSAGE: The NEF-shim is a thin FastAPI service that translates CAMARA calls into k8s API calls. Nothing else.\n\n"
        "Talking points:\n"
        "- Stateless control-plane component, sits in the LTH cluster.\n"
        "- Holds two kube-clients: in-cluster for LTH, kubeconfig file for Xerces.\n"
        "- Tenancy is by source-IP hash. No tokens, no signup — bring an IP, get a slug.\n"
        "- The red line at the bottom is THE most important point of the deck: inference traffic doesn't traverse this service.\n\n"
        "Colour meaning: Blue = external API, Teal = internal k8s operations, Coral = caveat / important note.")


def slide_tenant_isolation():
    slide = prs.slides.add_slide(prs.slide_layouts[4])
    set_title(slide, "Multi-tenancy by source IP")

    y = Inches(2.0)
    # Three tenant lanes
    tenants = [
        ("UE A   (5G IP 10.45.0.7)", "t-94a3d188", BLUE),
        ("UE B   (5G IP 10.45.0.8)", "t-13df96f6", TEAL),
        ("Laptop (corporate VPN)",  "t-660b0222", PURPLE),
    ]
    for i, (label, slug, color) in enumerate(tenants):
        ly = y + Inches(i * 1.4)
        add_box(slide, Inches(0.5), ly, Inches(4.0), Inches(0.9),
                label, fill_color=WHITE, border_color=color,
                font_size=12, text_color=color)
        add_arrow(slide, Inches(4.5), ly + Inches(0.45),
                  Inches(5.8), ly + Inches(0.45), color=color, width=Pt(2))
        add_textbox(slide, Inches(4.5), ly + Inches(0.05), Inches(1.3), Inches(0.3),
                    "md5(IP)[:8]", font_size=9, italic=True, color=GRAY2,
                    alignment=PP_ALIGN.CENTER)
        add_box(slide, Inches(5.8), ly, Inches(2.4), Inches(0.9),
                slug, fill_color=LIGHT_BLUE if color == BLUE else
                                  (LIGHT_TEAL if color == TEAL else LIGHT_PURPLE),
                border_color=color, font_size=14, bold=True, text_color=color,
                font_name=CODE_FONT)
        add_arrow(slide, Inches(8.3), ly + Inches(0.45),
                  Inches(9.0), ly + Inches(0.45), color=color, width=Pt(2))
        add_box(slide, Inches(9.0), ly, Inches(3.7), Inches(0.9),
                f"Deployment {slug}-app\nService     {slug}-app",
                fill_color=WHITE, border_color=color,
                font_size=11, font_name=CODE_FONT, text_color=GRAY1)

    add_textbox(slide, Inches(0.5), Inches(6.4), Inches(12.3), Inches(0.7),
                "One slug per source IP. Different IP → different slug → different pod and Service. "
                "No auth needed for the demo. VPN reconnect = new slug = orphan resources to clean.",
                font_size=12, color=GRAY1, italic=True)

    add_speaker_notes(slide,
        "KEY MESSAGE: Multi-tenancy is cheap — md5 of source IP and a k8s namespace label.\n\n"
        "Talking points:\n"
        "- Every CAMARA call carries x-forwarded-for. We hash, prefix t-, take 8 hex chars.\n"
        "- Each tenant gets its own Deployment + NodePort Service. Names collide deterministically.\n"
        "- Reality bite: switching network (5G → VPN → WiFi) changes the source IP. The old slug's resources become orphans because there's no /list endpoint to find them again.\n"
        "- In production this would be replaced with real auth tokens scoped to a customer.")


def slide_zones():
    slide = prs.slides.add_slide(prs.slide_layouts[4])
    set_title(slide, "Two zones — same API, same lifecycle")

    y = Inches(2.0)
    zones = [
        ("lth-5glab-gpu-zone",
         ["Provider: lth-kubernetes (in-house)",
          "GPU: 2 × NVIDIA L40S  (SM 89, Ada)",
          "CPU: 128 cores, 768 GB RAM",
          "Ingress: Envoy Gateway",
          "Reachable on: camara.5glab.control.lth.se"],
         BLUE),
        ("xerces-cloud-zone",
         ["Provider: ericsson-openstack (xerces)",
          "GPU: 1 × NVIDIA Tesla V100 (SM 70, Volta)",
          "CPU: 4 cores, 28 GB RAM",
          "Ingress: kube-proxy + OpenStack FIP",
          "Reachable on: 129.192.83.16"],
         TEAL),
    ]
    cw, gap = Inches(5.9), Inches(0.5)
    for i, (name, bullets, c) in enumerate(zones):
        x = Inches(0.5) + i * (cw + gap)
        add_box(slide, x, y, cw, Inches(0.8), name,
                fill_color=LIGHT_BLUE if c == BLUE else LIGHT_TEAL,
                border_color=c, font_size=16, bold=True, text_color=c,
                font_name=CODE_FONT)
        add_textbox(slide, x + Inches(0.2), y + Inches(1.0), cw - Inches(0.4),
                    Inches(3.0),
                    "\n".join("• " + b for b in bullets),
                    font_size=12, color=GRAY1, line_spacing=1.5)

    add_textbox(slide, Inches(0.5), Inches(6.0), Inches(12.3), Inches(0.6),
                "Client picks one based on capabilities. The endpoint discovery returns host:port "
                "with the right public host per zone — same shape, no proxy.",
                font_size=13, italic=True, color=GRAY1)

    add_speaker_notes(slide,
        "KEY MESSAGE: Two zones today, one API.\n\n"
        "Talking points:\n"
        "- LTH = our own k8s cluster on bare metal next to the gNB. Low latency.\n"
        "- Xerces = Ericsson's OpenStack cloud. More distant, older GPU.\n"
        "- The client compares capabilities and picks. The platform doesn't care.\n"
        "- The 'same shape' line is the punch — endpoint URLs look identical (host:port), only the host differs per zone.")


# ═══════════════════════════════════════════════════════════════════════════
#  PART 2 — THE APPLICATION
# ═══════════════════════════════════════════════════════════════════════════

def part2_divider():
    add_section_divider("Part 2 — EdgeVision",
                        "An application that maps onto the platform")


def slide_edgevision_intro():
    slide = prs.slides.add_slide(prs.slide_layouts[4])
    set_title(slide, "EdgeVision — what it is")

    y = Inches(2.0)
    add_textbox(slide, Inches(0.5), y, Inches(12.3), Inches(0.6),
                "Real-time YOLO object detection / segmentation, offloadable to the edge.",
                font_size=18, color=BLUE, bold=True)

    cards = [
        ("On the UE", [
            "Webcam-like frame source (camera-api)",
            "Tkinter GUI",
            "Local CPU detector (ultralytics)",
            "Remote detector (gRPC | HTTP)",
            "CAMARA client",
        ], BLUE),
        ("On the edge (per tenant)", [
            "Triton Inference Server",
            "Sidecar (decode → SHM → infer → NMS)",
            "TensorRT FP16 engines (per-GPU)",
            "Two models: YOLOv8n + YOLOv8x-seg",
            "gRPC :50051 + HTTP :8080 endpoints",
        ], TEAL),
    ]
    cw, gap = Inches(5.9), Inches(0.5)
    for i, (h, bullets, c) in enumerate(cards):
        x = Inches(0.5) + i * (cw + gap)
        add_box(slide, x, y + Inches(0.9), cw, Inches(0.7), h,
                fill_color=LIGHT_BLUE if c == BLUE else LIGHT_TEAL,
                border_color=c, font_size=15, bold=True, text_color=c)
        add_textbox(slide, x + Inches(0.2), y + Inches(1.7), cw - Inches(0.4),
                    Inches(3.0),
                    "\n".join("• " + b for b in bullets),
                    font_size=13, color=GRAY1, line_spacing=1.5)

    add_textbox(slide, Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.5),
                "Toggleable at runtime: E = LOCAL ↔ EDGE, T = gRPC ↔ HTTP, M = detect ↔ seg.",
                font_size=13, italic=True, color=GRAY2)

    add_speaker_notes(slide,
        "KEY MESSAGE: EdgeVision = one application split between a UE-side client and an edge-side inference pod.\n\n"
        "Talking points:\n"
        "- Anchor the audience: this is YOLO running on a phone vs running on a 5G-edge GPU.\n"
        "- The 'toggleable' bit is the demo's superpower: flip a key, watch latency change.\n"
        "- Both halves are stock open-source tech composed together.")


def slide_app_on_platform():
    """Same diagram as the generic platform slide — with EdgeVision boxes overlaid."""
    slide = prs.slides.add_slide(prs.slide_layouts[4])
    set_title(slide, "EdgeVision mapped onto the platform")
    add_textbox(slide, Inches(0.5), Inches(1.25), Inches(12.3), Inches(0.4),
                "Generic boxes from Part 1, now filled in with EdgeVision components.",
                font_size=14, italic=True, color=GRAY2)

    y_start = Inches(1.9)
    layers = [
        ("Client / UE", Inches(0.5), y_start, Inches(3.5), Inches(5.0), LIGHT_BLUE),
        ("5G SA Network", Inches(4.3), y_start, Inches(4.0), Inches(5.0), GRAY5),
        ("Edge Cloud (K8s)", Inches(8.6), y_start, Inches(4.2), Inches(5.0), LIGHT_TEAL),
    ]
    for label, lx, ly, lw, lh, fill in layers:
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, lx, ly, lw, lh)
        bg.fill.solid(); bg.fill.fore_color.rgb = fill
        bg.line.color.rgb = GRAY4; bg.line.width = Pt(0.5)
        _kill_shadow(bg)
        add_textbox(slide, lx + Inches(0.1), ly + Inches(0.08), lw - Inches(0.2), Inches(0.4),
                    label, font_size=12, bold=True, color=GRAY2)

    # Client side — concrete components
    cx = Inches(0.8)
    add_box(slide, cx, y_start + Inches(0.7), Inches(2.8), Inches(0.55),
            "GUI (Tkinter)", fill_color=WHITE, border_color=BLUE,
            font_size=12, bold=True, text_color=BLUE)
    add_box(slide, cx, y_start + Inches(1.35), Inches(2.8), Inches(0.55),
            "Local detector (CPU)", fill_color=WHITE, border_color=GRAY3,
            font_size=11, text_color=GRAY1)
    add_box(slide, cx, y_start + Inches(2.0), Inches(2.8), Inches(0.55),
            "Remote detector (gRPC | HTTP)", fill_color=WHITE, border_color=BLUE,
            font_size=11, text_color=BLUE)
    add_box(slide, cx, y_start + Inches(2.65), Inches(2.8), Inches(0.55),
            "CAMARA client", fill_color=LIGHT_BLUE, border_color=BLUE,
            font_size=11, bold=True, text_color=BLUE)
    add_box(slide, cx, y_start + Inches(3.3), Inches(2.8), Inches(0.55),
            "camera-api (MJPEG-ish)", fill_color=WHITE, border_color=GRAY3,
            font_size=11, text_color=GRAY1)

    # 5G core + NEF-shim
    nx = Inches(4.6)
    add_box(slide, nx, y_start + Inches(0.9), Inches(3.4), Inches(0.55),
            "gNB BB6651 + AIR6419", fill_color=WHITE, border_color=GRAY2,
            font_size=11, text_color=GRAY1)
    add_box(slide, nx, y_start + Inches(1.55), Inches(3.4), Inches(0.55),
            "Open5GS Core (AMF/SMF/UPF)", fill_color=WHITE, border_color=GRAY2,
            font_size=11, text_color=GRAY1)
    add_box(slide, nx, y_start + Inches(2.7), Inches(3.4), Inches(0.9),
            "NEF-shim\n(CAMARA control plane)",
            fill_color=LIGHT_BLUE, border_color=BLUE,
            font_size=12, bold=True, text_color=BLUE)

    # Edge side — concrete components
    ex = Inches(8.9)
    add_box(slide, ex, y_start + Inches(0.7), Inches(3.6), Inches(0.55),
            "Envoy Gateway", fill_color=WHITE, border_color=TEAL,
            font_size=11, text_color=TEAL)
    add_box(slide, ex, y_start + Inches(1.35), Inches(3.6), Inches(0.55),
            "NodePort  http :8080 / grpc :50051",
            fill_color=WHITE, border_color=TEAL,
            font_size=11, text_color=TEAL, font_name=CODE_FONT)
    add_box(slide, ex, y_start + Inches(2.0), Inches(3.6), Inches(0.7),
            "Sidecar (gRPC + HTTP, Python)",
            fill_color=WHITE, border_color=TEAL,
            font_size=11, bold=True, text_color=TEAL)
    add_box(slide, ex, y_start + Inches(2.8), Inches(3.6), Inches(0.7),
            "Triton  (yolov8n + yolov8x-seg)",
            fill_color=WHITE, border_color=TEAL,
            font_size=11, bold=True, text_color=TEAL)
    add_box(slide, ex, y_start + Inches(3.6), Inches(3.6), Inches(0.6),
            "NVIDIA L40S / V100  +  TensorRT FP16",
            fill_color=LIGHT_TEAL, border_color=TEAL,
            font_size=11, bold=True, text_color=TEAL)

    # Arrows
    # CAMARA control
    add_arrow(slide, Inches(3.6), y_start + Inches(2.95), Inches(4.6), y_start + Inches(3.0),
              color=BLUE, width=Pt(2))
    add_textbox(slide, Inches(3.55), y_start + Inches(3.05), Inches(1.4), Inches(0.3),
                "CAMARA", font_size=10, italic=True, color=BLUE)
    # NEF → k8s (control)
    add_arrow(slide, Inches(8.0), y_start + Inches(3.1), Inches(8.9), y_start + Inches(3.0),
              color=GRAY2, width=Pt(1.5))
    add_textbox(slide, Inches(8.0), y_start + Inches(3.2), Inches(1.0), Inches(0.3),
                "k8s API", font_size=10, italic=True, color=GRAY2)
    # Data plane: direct
    add_arrow(slide, Inches(3.6), y_start + Inches(2.3), Inches(8.9), y_start + Inches(2.3),
              color=TEAL, width=Pt(2))
    add_textbox(slide, Inches(5.0), y_start + Inches(1.95), Inches(3.2), Inches(0.3),
                "inference traffic (direct to NodePort)", font_size=11,
                italic=True, color=TEAL, alignment=PP_ALIGN.CENTER)

    add_speaker_notes(slide,
        "KEY MESSAGE: Same three columns as Part 1 — now with EdgeVision parts in each box.\n\n"
        "Talking points:\n"
        "- Side-by-side this with the Part 1 diagram. The platform didn't change.\n"
        "- The Tkinter GUI is just a CAMARA client + a gRPC/HTTP client; nothing CAMARA-specific is hardcoded.\n"
        "- The edge pod is whatever the manifest says — here it's Triton + sidecar. Same NodePorts as any other pod.\n"
        "- The teal arrow goes straight from client to pod. CAMARA isn't on the hot path.\n\n"
        "Colour meaning: same as Part 1 (Blue control, Teal data plane).")


def slide_app_manifest():
    slide = prs.slides.add_slide(prs.slide_layouts[4])
    set_title(slide, "Application manifest")

    y = Inches(1.9)
    # Code block on the left
    manifest = (
        "APP_MANIFEST = {\n"
        "  \"appName\":             \"edgevision-yolov8\",\n"
        "  \"appProvider\":         \"lth-frtn90\",\n"
        "  \"appSoftwareVersion\":  \"1.0.0\",\n"
        "  \"packageType\":         \"CONTAINER\",\n"
        "  \"containerSpec\": {\n"
        "    \"imageRegistry\": \"ghcr.io/ekeroid/5glab\",\n"
        "    \"imageName\":     \"edgevision-infer\",\n"
        "    \"imageTag\":      \"latest\",\n"
        "    \"readinessProbe\": {\n"
        "      \"http\": {\"path\": \"/health\", \"port\": 8080},\n"
        "      \"initialDelaySeconds\": 30, \"periodSeconds\": 5},\n"
        "    \"livenessProbe\": {\n"
        "      \"http\": {\"path\": \"/health\", \"port\": 8080},\n"
        "      \"initialDelaySeconds\": 1200,\n"
        "      \"periodSeconds\": 15, \"failureThreshold\": 5}\n"
        "  },\n"
        "  \"requiredResources\": {\n"
        "    \"cpu\": 2, \"memory\": 8192, \"gpu\": 1\n"
        "  },\n"
        "  \"componentSpec\": [{\n"
        "    \"componentName\": \"infer\",\n"
        "    \"networkInterfaces\": [\n"
        "      {\"name\": \"http\",  \"port\": 8080,  \"protocol\": \"TCP\"},\n"
        "      {\"name\": \"grpc\",  \"port\": 50051, \"protocol\": \"TCP\"}\n"
        "    ]\n"
        "  }]\n"
        "}\n"
    )
    add_box(slide, Inches(0.5), y, Inches(7.4), Inches(5.0), "",
            fill_color=GRAY5, border_color=GRAY3, rounded=False)
    add_textbox(slide, Inches(0.7), y + Inches(0.15), Inches(7.0), Inches(4.8),
                manifest, font_size=10, font_name=CODE_FONT, color=GRAY1,
                line_spacing=1.15)

    # Right-side annotations
    ax = Inches(8.2)
    add_textbox(slide, ax, y, Inches(4.7), Inches(0.4),
                "POST /edge-app-management/v0/apps",
                font_size=12, bold=True, color=BLUE, font_name=CODE_FONT)
    annotations = [
        ("imageTag", "Same image runs on both zones."),
        ("livenessProbe 1200s",
         "Wide grace — V100 cold builds TRT engine in ~12 min."),
        ("gpu: 1",
         "Sets runtimeClassName=nvidia, GPU resource request, GPU taint toleration."),
        ("componentSpec ports",
         "Both forwarded as NodePort. Name passed through to endpoint discovery."),
    ]
    ty = y + Inches(0.6)
    for label, body in annotations:
        add_textbox(slide, ax, ty, Inches(4.7), Inches(0.35),
                    label, font_size=12, bold=True, color=TEAL,
                    font_name=CODE_FONT)
        add_textbox(slide, ax, ty + Inches(0.3), Inches(4.7), Inches(0.8),
                    body, font_size=11, color=GRAY1)
        ty += Inches(1.1)

    add_speaker_notes(slide,
        "KEY MESSAGE: The whole 'application' is one JSON document. The platform does the rest.\n\n"
        "Talking points:\n"
        "- Container image is fixed; lifecycle is per-instance.\n"
        "- The two NICs (http, grpc) become two NodePort ports. Names are preserved in endpoint discovery.\n"
        "- 1200s liveness grace is calibrated for the V100 worst case. On L40S we never come close to that.\n"
        "- gpu=1 is the trigger for GPU scheduling, taint toleration, runtimeClass — all done by the shim.\n\n"
        "REFERENCE: client/camara.py — this is the literal APP_MANIFEST dict sent on POST /apps.")


def slide_container_contents():
    slide = prs.slides.add_slide(prs.slide_layouts[4])
    set_title(slide, "What's inside edgevision-infer:latest")

    y = Inches(2.0)
    rows = [
        ("Base",         "nvcr.io/nvidia/tritonserver:23.10-py3",
         "TensorRT 8.6 — supports SM 70 (Volta/V100) through SM 89 (Ada/L40S)."),
        ("Runtime",      "fastapi · uvicorn · tritonclient[grpc] · opencv · ultralytics",
         "Sidecar deps + export toolchain (ONNX, onnxslim)."),
        ("Models",       "yolov8n.pt   yolov8x-seg.pt",
         "Ultralytics weights. .plan engines are NOT baked — built at first boot per GPU."),
        ("Model repo",   "models/yolov8n_trt/config.pbtxt  yolov8xseg_trt/config.pbtxt",
         "Triton model repository skeleton."),
        ("Sidecar",      "sidecar/main.py  +  infer_pb2*.py",
         "gRPC :50051 (InferenceService) + HTTP :8080 (/infer, /health)."),
        ("Tooling",      "export_model.py  entrypoint.sh",
         ".pt → ONNX → TRT FP16 plan. Entrypoint orchestrates build + Triton + sidecar."),
    ]
    col_w = [Inches(1.6), Inches(5.0), Inches(6.2)]
    col_x = [Inches(0.5), Inches(2.2), Inches(7.3)]
    headers = ["Section", "Files", "Notes"]
    for j, (h, cx, cw) in enumerate(zip(headers, col_x, col_w)):
        add_box(slide, cx, y, cw, Inches(0.45), h,
                fill_color=BLUE, border_color=BLUE, font_size=12,
                bold=True, text_color=WHITE, rounded=False)
    for i, (section, files, notes) in enumerate(rows):
        ry = y + Inches(0.5 + i * 0.65)
        fill = WHITE if i % 2 == 0 else GRAY5
        add_box(slide, col_x[0], ry, col_w[0], Inches(0.6), section,
                fill_color=fill, border_color=GRAY4, font_size=11,
                bold=True, text_color=GRAY1, rounded=False)
        add_box(slide, col_x[1], ry, col_w[1], Inches(0.6), files,
                fill_color=fill, border_color=GRAY4, font_size=10,
                text_color=GRAY1, rounded=False, font_name=CODE_FONT)
        add_box(slide, col_x[2], ry, col_w[2], Inches(0.6), notes,
                fill_color=fill, border_color=GRAY4, font_size=10,
                text_color=GRAY1, rounded=False)

    add_textbox(slide, Inches(0.5), Inches(6.4), Inches(12.3), Inches(0.6),
                "Engines are GPU-arch-specific → cached in pod-scoped emptyDir → "
                "first boot rebuilds on every fresh pod.",
                font_size=12, italic=True, color=CORAL)

    add_speaker_notes(slide,
        "KEY MESSAGE: Image is just Triton + sidecar + weights + a build script. No pre-built engines.\n\n"
        "Talking points:\n"
        "- We rebased to 23.10 specifically so the same image works on V100. The earlier 24.10 base dropped SM 70.\n"
        "- No .plan files in the image means the same artifact is portable across GPU generations.\n"
        "- The cost: first boot per fresh pod. L40S ~3-5 min, V100 ~12-14 min.\n"
        "- A PersistentVolume on the engine cache directory would fix this. Future work.")


def slide_first_boot():
    slide = prs.slides.add_slide(prs.slide_layouts[4])
    set_title(slide, "First boot — what the entrypoint does")

    y = Inches(2.0)
    steps = [
        ("1", "Export .pt → ONNX  (yolov8n)",                  "~10–30 s on L40S, ~5 min on V100"),
        ("2", "trtexec → TensorRT FP16 plan (yolov8n)",        "Same range"),
        ("3", "Export .pt → ONNX  (yolov8x-seg)",              "Bigger model"),
        ("4", "trtexec → TensorRT FP16 plan (yolov8x-seg)",    "~3-5 min L40S, ~6-8 min V100"),
        ("5", "Start Triton  (load both models)",              "~10 s"),
        ("6", "Start sidecar (gRPC :50051 + HTTP :8080)",      "Instant"),
        ("7", "Readiness probe passes → NodePort accepts traffic", "Endpoint discovery now valid"),
    ]
    for i, (n, what, when) in enumerate(steps):
        ly = y + Inches(i * 0.55)
        add_box(slide, Inches(0.5), ly, Inches(0.5), Inches(0.45), n,
                fill_color=BLUE, border_color=BLUE, font_size=12, bold=True,
                text_color=WHITE)
        add_textbox(slide, Inches(1.1), ly, Inches(7.0), Inches(0.45),
                    what, font_size=12, color=GRAY1, font_name=CODE_FONT)
        add_textbox(slide, Inches(8.4), ly, Inches(4.4), Inches(0.45),
                    when, font_size=11, italic=True, color=GRAY2)

    # Summary timings
    sy = y + Inches(4.2)
    add_textbox(slide, Inches(0.5), sy, Inches(12.3), Inches(0.4),
                "Total cold start:",
                font_size=14, bold=True, color=BLACK)
    add_box(slide, Inches(0.5), sy + Inches(0.5), Inches(5.9), Inches(0.7),
            "L40S:  3 – 5  min",
            fill_color=LIGHT_TEAL, border_color=TEAL,
            font_size=18, bold=True, text_color=TEAL)
    add_box(slide, Inches(6.9), sy + Inches(0.5), Inches(5.9), Inches(0.7),
            "V100:  12 – 14  min",
            fill_color=LIGHT_CORAL, border_color=CORAL,
            font_size=18, bold=True, text_color=CORAL)

    add_speaker_notes(slide,
        "KEY MESSAGE: First boot is dominated by TRT engine build, and is GPU-dependent.\n\n"
        "Talking points:\n"
        "- Important to call out the V100 cold time, otherwise the liveness probe slide doesn't make sense.\n"
        "- A warm restart (same node, same pod's emptyDir) skips steps 1-4. ~10 s to ready.\n"
        "- A new pod on the same node STILL rebuilds because emptyDir is pod-scoped. Future: move engines to a PVC.\n\n"
        "Colour meaning: Teal = good (L40S), Coral = caveat (V100).")


def slide_deploy_sequence():
    slide = prs.slides.add_slide(prs.slide_layouts[4])
    set_title(slide, "Deployment sequence")

    y = Inches(1.9)
    actors = [
        ("Client",   Inches(1.3)),
        ("NEF-shim", Inches(4.5)),
        ("K8s API",  Inches(7.7)),
        ("Pod",      Inches(10.7)),
    ]
    for name, x in actors:
        c = LIGHT_BLUE if name in ("Client", "NEF-shim") else LIGHT_TEAL
        b = BLUE if name in ("Client", "NEF-shim") else TEAL
        t = BLUE if name in ("Client", "NEF-shim") else TEAL
        add_box(slide, x - Inches(0.6), y, Inches(1.4), Inches(0.5),
                name, fill_color=c, border_color=b,
                font_size=12, bold=True, text_color=t)
        add_line(slide, x + Inches(0.1), y + Inches(0.5),
                 x + Inches(0.1), y + Inches(4.7),
                 color=GRAY3, width=Pt(1), dashed=True)

    messages = [
        (0, 1, "GET /edge-cloud-zones",            Inches(0.7)),
        (1, 0, "zones[]",                          Inches(0.95)),
        (0, 1, "POST /apps  (manifest)",           Inches(1.2)),
        (1, 0, "appId",                            Inches(1.45)),
        (0, 1, "POST /app-instances",              Inches(1.7)),
        (1, 2, "create Deployment + Service",      Inches(1.95)),
        (2, 3, "schedule + image pull",            Inches(2.2)),
        (3, 3, "build TRT engines  (3–14 min)",    Inches(2.5)),
        (0, 1, "GET /app-instances/{id}  (poll)",  Inches(2.85)),
        (1, 2, "watch pod status",                 Inches(3.1)),
        (3, 2, "readinessProbe passes",            Inches(3.35)),
        (1, 0, "status: ready",                    Inches(3.6)),
        (0, 1, "GET /endpoints",                   Inches(3.85)),
        (1, 0, "{http: host:port, grpc: host:port}", Inches(4.1)),
    ]
    for src, dst, label, dy in messages:
        sx = actors[src][1] + Inches(0.1)
        dx = actors[dst][1] + Inches(0.1)
        msg_y = y + dy
        if src == dst:
            add_arrow(slide, sx + Inches(0.3), msg_y, sx + Inches(0.3),
                      msg_y + Inches(0.3), color=TEAL, width=Pt(1.5))
            add_textbox(slide, sx + Inches(0.4), msg_y - Inches(0.05),
                        Inches(2.8), Inches(0.25), label,
                        font_size=10, italic=True, color=TEAL)
        else:
            add_arrow(slide, sx, msg_y, dx, msg_y, color=BLUE, width=Pt(1.5))
            mid_x = min(sx, dx) + Inches(0.1)
            add_textbox(slide, mid_x, msg_y - Inches(0.25),
                        Inches(3.0), Inches(0.25), label,
                        font_size=10, color=GRAY1)

    add_speaker_notes(slide,
        "KEY MESSAGE: One toggle (E in the GUI) drives a 5-call CAMARA sequence and a few minutes of k8s + Triton startup.\n\n"
        "Talking points:\n"
        "- The polling loop is the GUI showing real-time progress (scheduling → pulling → starting → running → ready).\n"
        "- Engine build is the dominant cost first time. After that the same pod hits the cache.\n"
        "- The last reply ('endpoints') is the entire reason CAMARA is on the slide. Hand-off complete; from here it's direct gRPC/HTTP.")


def slide_inference_sequence():
    slide = prs.slides.add_slide(prs.slide_layouts[4])
    set_title(slide, "Inference sequence (steady state)")

    y = Inches(1.9)
    actors = [
        ("Client",  Inches(1.2)),
        ("Sidecar", Inches(4.5)),
        ("Triton",  Inches(7.7)),
        ("GPU",     Inches(10.7)),
    ]
    for name, x in actors:
        c = LIGHT_BLUE if "Client" in name else LIGHT_TEAL
        b = BLUE if "Client" in name else TEAL
        add_box(slide, x - Inches(0.55), y, Inches(1.3), Inches(0.5),
                name, fill_color=c, border_color=b,
                font_size=12, bold=True, text_color=b)
        add_line(slide, x + Inches(0.1), y + Inches(0.5),
                 x + Inches(0.1), y + Inches(4.7),
                 color=GRAY3, width=Pt(1), dashed=True)

    messages = [
        (0, 1, "gRPC: InferRequest{jpeg, model}",  Inches(0.7)),
        (1, 1, "decode + resize 640×640",           Inches(1.1)),
        (1, 2, "SHM write + infer()",               Inches(1.5)),
        (2, 3, "TensorRT FP16 forward",             Inches(1.9)),
        (3, 2, "output tensors",                    Inches(2.3)),
        (2, 1, "SHM read",                          Inches(2.7)),
        (1, 1, "NMS / mask decode",                 Inches(3.1)),
        (1, 0, "InferResponse{detections[]}",       Inches(3.5)),
    ]
    for src, dst, label, dy in messages:
        sx = actors[src][1] + Inches(0.1)
        dx = actors[dst][1] + Inches(0.1)
        msg_y = y + dy
        if src == dst:
            add_arrow(slide, sx + Inches(0.2), msg_y, sx + Inches(0.2),
                      msg_y + Inches(0.25), color=TEAL, width=Pt(1.5))
            add_textbox(slide, sx + Inches(0.3), msg_y - Inches(0.05),
                        Inches(2.8), Inches(0.25), label,
                        font_size=10, italic=True, color=TEAL)
        else:
            add_arrow(slide, sx, msg_y, dx, msg_y, color=BLUE, width=Pt(1.5))
            mid_x = min(sx, dx) + Inches(0.1)
            add_textbox(slide, mid_x, msg_y - Inches(0.25),
                        Inches(3.2), Inches(0.25), label,
                        font_size=10, color=GRAY1)

    add_textbox(slide, Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.5),
                "Client ↔ Pod direct. NEF-shim is not in this picture.",
                font_size=14, bold=True, color=CORAL, alignment=PP_ALIGN.CENTER)

    add_speaker_notes(slide,
        "KEY MESSAGE: Inference is a 4-actor sequence inside one pod plus a network hop.\n\n"
        "Talking points:\n"
        "- Sidecar ↔ Triton is in-process: shared-memory IPC, no socket.\n"
        "- The visible network hop is Client ↔ Sidecar only. NodePort, single TCP for gRPC.\n"
        "- The bottom red note is the headline — repeat it.")


def slide_models():
    slide = prs.slides.add_slide(prs.slide_layouts[4])
    set_title(slide, "Two models — switchable at runtime")

    y = Inches(2.0)
    cards = [
        ("YOLOv8n — Object Detection", BLUE,
         ["6.3 M parameters",
          "Input: 640×640×3 (FP16)",
          "Output: [1, 84, 8400] (boxes + 80 classes)",
          "Edge L40S: ~4 ms per frame",
          "Local CPU (M1 Pro): ~120 ms",
          "Use: counting, tracking, presence."]),
        ("YOLOv8x-seg — Instance Segmentation", BLUE,
         ["71.8 M parameters",
          "Input: 640×640×3 (FP16)",
          "Output 0: [1, 116, 8400] (boxes + mask coeffs)",
          "Output 1: [1, 32, 160, 160] (prototypes)",
          "Edge L40S: ~6 ms per frame",
          "Local CPU (M1 Pro): ~530 ms",
          "Use: pixel-level boundaries, AR overlays."]),
    ]
    cw, gap = Inches(5.9), Inches(0.5)
    for i, (h, c, bullets) in enumerate(cards):
        x = Inches(0.5) + i * (cw + gap)
        add_box(slide, x, y, cw, Inches(0.7), h,
                fill_color=LIGHT_BLUE, border_color=c,
                font_size=16, bold=True, text_color=c)
        add_textbox(slide, x + Inches(0.2), y + Inches(0.85), cw - Inches(0.4),
                    Inches(4.0),
                    "\n".join("• " + b for b in bullets),
                    font_size=13, color=GRAY1, line_spacing=1.5)

    add_textbox(slide, Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.5),
                "Both models live in the same Triton model repo. Switch by changing one header.",
                font_size=13, italic=True, color=GRAY2)

    add_speaker_notes(slide,
        "KEY MESSAGE: Two YOLOv8 models, same image, runtime switch.\n\n"
        "Talking points:\n"
        "- Detection is the cheap one — even local CPU can keep up at low frame rates.\n"
        "- Segmentation is the gain case — 88× speedup vs local CPU.\n"
        "- 'Same Triton model repo' is the practical detail — no second pod to manage.")


def slide_runtime_latency():
    slide = prs.slides.add_slide(prs.slide_layouts[4])
    set_title(slide, "Measured latency — gRPC over 5G SA, LTH zone")

    y = Inches(1.95)
    # Two side-by-side tables
    def render_table(title_text, rows, x, headercolor):
        add_textbox(slide, x, y, Inches(6.2), Inches(0.4),
                    title_text, font_size=14, bold=True, color=headercolor)
        col_x = [x, x + Inches(2.2), x + Inches(3.5)]
        col_w = [Inches(2.2), Inches(1.3), Inches(1.2)]
        headers = ["Segment", "Avg", "P95"]
        for j, (h, cx, cw) in enumerate(zip(headers, col_x, col_w)):
            add_box(slide, cx, y + Inches(0.5), cw, Inches(0.4), h,
                    fill_color=headercolor, border_color=headercolor,
                    font_size=11, bold=True, text_color=WHITE, rounded=False)
        for i, (seg, avg, p95) in enumerate(rows):
            ry = y + Inches(0.9 + i * 0.42)
            fill = WHITE if i % 2 == 0 else GRAY5
            for j, (txt, cx, cw) in enumerate(zip([seg, avg, p95], col_x, col_w)):
                add_box(slide, cx, ry, cw, Inches(0.4), txt,
                        fill_color=fill, border_color=GRAY4, font_size=10,
                        text_color=GRAY1, rounded=False,
                        font_name=CODE_FONT if j > 0 else BODY_FONT)

    detect_rows = [
        ("End-to-end",      "55.0 ms",  "73.1 ms"),
        ("5G UL + DL",      "31.5 ms",  "45.9 ms"),
        ("Server total",    "23.5 ms",  "26.7 ms"),
        ("  Preprocess",    "16.9 ms",  "19.6 ms"),
        ("  GPU inference", " 3.9 ms",  " 4.2 ms"),
        ("  Postprocess",   " 2.7 ms",  " 3.3 ms"),
    ]
    seg_rows = [
        ("End-to-end",      "50.7 ms",  "66.6 ms"),
        ("5G UL + DL",      "31.3 ms",  "42.6 ms"),
        ("Server total",    "19.3 ms",  "30.1 ms"),
        ("  Preprocess",    " 8.9 ms",  "16.7 ms"),
        ("  GPU inference", " 5.6 ms",  " 6.0 ms"),
        ("  Postprocess",   " 4.7 ms",  " 7.6 ms"),
    ]
    render_table("Detection — YOLOv8n (195 samples)",      detect_rows, Inches(0.5), BLUE)
    render_table("Segmentation — YOLOv8x-seg (200 samples)", seg_rows,   Inches(6.85), TEAL)

    # Local CPU comparison footer
    fy = Inches(6.0)
    add_textbox(slide, Inches(0.5), fy, Inches(12.3), Inches(0.4),
                "Local CPU baseline (M1 Pro, same client):",
                font_size=12, bold=True, color=GRAY1)
    add_textbox(slide, Inches(0.5), fy + Inches(0.4), Inches(12.3), Inches(0.5),
                "  YOLOv8n      48.9 ms avg  /  54.1 ms P95            "
                "YOLOv8x-seg   344.9 ms avg  /  423.8 ms P95",
                font_size=12, font_name=CODE_FONT, color=GRAY1)

    add_speaker_notes(slide,
        "KEY MESSAGE: Real measurements from the running system, not estimates.\n\n"
        "Talking points:\n"
        "- End-to-end ≈ 50–55 ms over 5G. Of that, the radio is ~30 ms and the pod is ~20 ms.\n"
        "- GPU inference proper is 4-6 ms. Everything else is movement + decode + NMS.\n"
        "- Segmentation is barely slower than detection on edge GPU. On CPU it's 7× slower.\n"
        "- The CPU baseline at the bottom is to anchor why edge matters at all.")


def slide_transport():
    slide = prs.slides.add_slide(prs.slide_layouts[4])
    set_title(slide, "Transport choice — gRPC vs HTTP")

    y = Inches(2.0)
    cards = [
        ("gRPC  (primary)", BLUE, [
            "Direct to NodePort, host:port",
            "HTTP/2 binary framing, persistent connection",
            "InferenceService/Infer (proto: jpeg + model)",
            "~50 ms end-to-end over 5G (LTH zone)",
            "Reliable on 5G; blocked through corporate VPN proxies",
        ]),
        ("HTTP  (fallback)", TEAL, [
            "Direct to NodePort, http://host:port",
            "JPEG POST  /infer  →  JSON response",
            "New TCP connection per request",
            "~5–7 ms slower (encode + parse overhead)",
            "Works wherever TCP works. Survives proxies.",
        ]),
    ]
    cw, gap = Inches(5.9), Inches(0.5)
    for i, (h, c, bullets) in enumerate(cards):
        x = Inches(0.5) + i * (cw + gap)
        add_box(slide, x, y, cw, Inches(0.7), h,
                fill_color=LIGHT_BLUE if c == BLUE else LIGHT_TEAL,
                border_color=c, font_size=16, bold=True, text_color=c)
        add_textbox(slide, x + Inches(0.2), y + Inches(0.85), cw - Inches(0.4),
                    Inches(3.5),
                    "\n".join("• " + b for b in bullets),
                    font_size=13, color=GRAY1, line_spacing=1.5)

    add_textbox(slide, Inches(0.5), Inches(6.0), Inches(12.3), Inches(0.6),
                "Same endpoint shape (host:port). Same NodePort. Choice is in the client.",
                font_size=13, italic=True, color=GRAY1)

    add_speaker_notes(slide,
        "KEY MESSAGE: Two transports, same NodePort. Client picks; platform doesn't care.\n\n"
        "Talking points:\n"
        "- gRPC = faster + persistent. HTTP = robust through proxies.\n"
        "- The corporate-VPN gotcha is real — was the reason we changed the endpoint discovery to return host:port for both.\n"
        "- 'Same endpoint shape' is the key — there is no special path-prefix proxy any more.")


# ═══════════════════════════════════════════════════════════════════════════
#  CLOSING
# ═══════════════════════════════════════════════════════════════════════════

def slide_closing():
    slide = prs.slides.add_slide(prs.slide_layouts[20])  # Black logo end slide
    subtitle_ph = slide.placeholders[1]
    subtitle_ph.text = "EdgeVision — LTH 5G SA Lab"
    for p in subtitle_ph.text_frame.paragraphs:
        for run in p.runs:
            run.font.name = BODY_FONT
            run.font.size = Pt(14)
            run.font.color.rgb = WHITE
    add_speaker_notes(slide,
        "KEY MESSAGE: Thanks. Questions welcome.\n\n"
        "Talking points:\n"
        "- Demo is live and can be shown.\n"
        "- Code: edgevision/ in the project repo.\n"
        "- Open follow-ups: PVC-backed engine cache, real auth, more zones, model hot-swap.\n")


# ═══════════════════════════════════════════════════════════════════════════
#  BUILD
# ═══════════════════════════════════════════════════════════════════════════

print("Building EdgeVision-Overview v2 …")

slide_title()

part1_divider()
slide_why_offload()
slide_what_is_camara()
slide_camara_state()
slide_platform_arch_generic()
slide_lifecycle_generic()
slide_camara_api_details()
slide_camara_sequence_generic()
slide_camara_curl()
slide_nef_shim_role()
slide_tenant_isolation()
slide_zones()

part2_divider()
slide_edgevision_intro()
slide_app_on_platform()
slide_app_manifest()
slide_container_contents()
slide_first_boot()
slide_deploy_sequence()
slide_inference_sequence()
slide_models()
slide_runtime_latency()
slide_transport()

slide_closing()


# ─── Slide numbers ────────────────────────────────────────────────────────────
for i, slide in enumerate(prs.slides):
    if i == 0 or i == len(prs.slides) - 1:
        continue
    txBox = slide.shapes.add_textbox(Inches(12.0), Inches(7.05), Inches(0.8), Inches(0.3))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = str(i + 1)
    p.alignment = PP_ALIGN.RIGHT
    run = p.runs[0]
    run.font.size = Pt(10)
    run.font.color.rgb = GRAY2
    run.font.name = BODY_FONT
    _kill_shadow(txBox)


# ─── Final pass: kill inherited shadows ───────────────────────────────────────
for slide in prs.slides:
    _kill_all_shadows_on_slide(slide)


prs.save(OUTPUT_PATH)
print(f"\nSaved: {OUTPUT_PATH}")
print(f"Slides: {len(prs.slides)}")
